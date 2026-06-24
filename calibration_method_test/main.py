import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches

# ==========================
# 1️⃣ Load ROI per file
# ==========================
def extract_roi(df):
    """Remove last 3 reference columns and return full 2D matrix of pixels > 1pF"""
    if df is None:
        return None
    df = df.iloc[:, :-3]
    matrix = df.apply(pd.to_numeric, errors='coerce').values
    mask = (~np.isnan(matrix)) & (matrix > 1.0)
    if not np.any(mask):
        return None
    r_min, r_max = np.where(np.any(mask, axis=1))[0][[0, -1]]
    c_min, c_max = np.where(np.any(mask, axis=0))[0][[0, -1]]
    roi = matrix[r_min:r_max+1, c_min:c_max+1]
    return roi

def load_calibration(folder, pressures):
    """Load all pressures into DataFrame with full 2D ROI"""
    data = []
    for p in pressures:
        up_path = os.path.join(folder, f"{p}u.csv")
        down_path = os.path.join(folder, f"{p}d.csv")
        if not os.path.exists(up_path):
            continue
        df_up = pd.read_csv(up_path, skiprows=1, index_col=0)
        df_down = pd.read_csv(down_path, skiprows=1, index_col=0) if os.path.exists(down_path) else None
        roi_up = extract_roi(df_up)
        roi_down = extract_roi(df_down) if df_down is not None else None
        if roi_up is None:
            continue
        data.append({
            "P": p,
            "C_up_vals": roi_up,
            "C_down_vals": roi_down
        })
    return pd.DataFrame(data)

# ==========================
# 2️⃣ Compute common ROI shape
# ==========================
def get_common_shape(df_list):
    """Find minimal rows and cols across all DataFrames"""
    min_rows, min_cols = np.inf, np.inf
    for df in df_list:
        for idx in range(len(df)):
            arr = df.iloc[idx]["C_up_vals"]
            if arr.ndim == 1:
                arr = arr.reshape(-1,1)
            r, c = arr.shape
            min_rows = min(min_rows, r)
            min_cols = min(min_cols, c)
    return int(min_rows), int(min_cols)

# ==========================
# 3️⃣ Flatten ROIs to common ROI
# ==========================
def flatten_common_roi(df, min_rows, min_cols):
    """Crop each ROI to common minimal shape and flatten"""
    for idx in range(len(df)):
        arr_up = df.iloc[idx]["C_up_vals"]
        arr_down = df.iloc[idx]["C_down_vals"] if df.iloc[idx]["C_down_vals"] is not None else None
        if arr_up.ndim == 1:
            arr_up = arr_up.reshape(-1,1)
        arr_up_crop = arr_up[:min_rows,:min_cols].flatten()
        arr_down_crop = arr_down[:min_rows,:min_cols].flatten() if arr_down is not None else None
        df.at[idx,"C_up_vals"] = arr_up_crop
        df.at[idx,"C_down_vals"] = arr_down_crop
    return df, min_rows*min_cols

# ==========================
# 4️⃣ Fit polynomial per pixel
# ==========================
def fit_poly_per_pixel(df, degree=3):
    P = df["P"].values
    num_pixels = len(df.iloc[0]["C_up_vals"])
    coeffs = np.zeros((num_pixels, degree+1))
    for i in range(num_pixels):
        C = np.array([df.iloc[j]["C_up_vals"][i] for j in range(len(df))])
        coeffs[i,:] = np.polyfit(P, C, degree)
    return coeffs

# ==========================
# 5️⃣ Create report
# ==========================
def create_pixel_polynomial_report(path70, path110, output="Pixel_Poly_Report.pptx"):
    pressures70 = [0,10,30,50,70]
    pressures110 = [0,10,30,50,70,90,110]

    df70 = load_calibration(path70, pressures70)
    df110 = load_calibration(path110, pressures110)

    # Find global common ROI
    min_rows, min_cols = get_common_shape([df70, df110])
    df70, n_pixels = flatten_common_roi(df70, min_rows, min_cols)
    df110, _ = flatten_common_roi(df110, min_rows, min_cols)

    # Fit polynomials
    coeffs70 = fit_poly_per_pixel(df70)
    coeffs110 = fit_poly_per_pixel(df110[df110.P<=70])

    # Compare coefficients
    coeff_diff = coeffs70 - coeffs110
    mean_diff = np.mean(np.abs(coeff_diff), axis=0)

    # RMSE per pixel for pressure prediction using poly70
    rmse_list = []
    for i in range(n_pixels):
        poly = np.poly1d(coeffs70[i])
        C_vals = np.array([df110.iloc[j]["C_up_vals"][i] for j in range(len(df110[df110.P<=70]))])
        P_real = df110[df110.P<=70]["P"].values
        P_pred = poly(C_vals)
        rmse_list.append(np.sqrt(np.mean((P_pred - P_real)**2)))
    rmse_pressure = np.mean(rmse_list)

    # ==========================
    # Visualization
    # ==========================
    fig, axes = plt.subplots(2,2, figsize=(16,12))

    # Mean Coeff Diff per term
    axes[0,0].bar(['a','b','c','d'], mean_diff)
    axes[0,0].set_title("Mean Coefficient Differences (70 vs 110)")
    axes[0,0].set_ylabel("Δ coefficient")

    # RMSE histogram per pixel
    axes[0,1].hist(rmse_list, bins=20, color='purple')
    axes[0,1].set_title("Pixel-wise Pressure RMSE (Model 70 on Data 110)")
    axes[0,1].set_xlabel("RMSE [Pressure Units]")

    # Example: first 5 pixel polynomials
    for i in range(min(5,n_pixels)):
        poly70 = np.poly1d(coeffs70[i])
        poly110 = np.poly1d(coeffs110[i])
        P_plot = np.linspace(0,70,100)
        axes[1,0].plot(P_plot, poly70(P_plot), 'b--', alpha=0.7)
        axes[1,0].plot(P_plot, poly110(P_plot), 'r-', alpha=0.7)
    axes[1,0].set_title("Example Pixel Polynomials (70 vs 110)")
    axes[1,0].set_xlabel("Pressure")
    axes[1,0].set_ylabel("Capacitance [pF]")

    # Mean polynomial difference across pixels
    axes[1,1].plot(np.arange(n_pixels), np.linalg.norm(coeff_diff, axis=1))
    axes[1,1].set_title("Polynomial Differences per Pixel")
    axes[1,1].set_xlabel("Pixel Index")
    axes[1,1].set_ylabel("Norm of Δ coefficients")

    plt.tight_layout()
    fig_path = "pixel_poly_analysis.png"
    plt.savefig(fig_path)
    plt.close(fig)

    # ==========================
    # PowerPoint
    # ==========================
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Pixel-wise Calibration Comparison"
    slide.placeholders[1].text = "ROI Minimum shared across all files; Cubic Polynomials per pixel"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Polynomial Analysis"
    slide.shapes.add_picture(fig_path, Inches(0.5), Inches(1.5), width=Inches(9))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Mean Coefficient Differences & RMSE"
    tf = slide.placeholders[1].text_frame
    tf.text = f"Mean Δ coefficients (a,b,c,d): {mean_diff}"
    tf.add_paragraph().text = f"Mean Pixel RMSE: {rmse_pressure:.4f}"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Final Recommendation"
    tf = slide.placeholders[1].text_frame
    if rmse_pressure < 1.5:
        tf.text = "Calibration up to 70 is sufficient."
        tf.add_paragraph().text = "Pixel-level models show minimal differences."
    else:
        tf.text = "Calibration up to 110 recommended."
        tf.add_paragraph().text = "Pixel-level models indicate changes after 70."

    prs.save(output)
    print(f"Report generated: {output}")

# ==========================
# Run
# ==========================
if __name__ == "__main__":
    path_70 = r"C:\Users\dvirs\Downloads\0-70\files"
    path_110 = r"C:\Users\dvirs\Downloads\0-110\files"
    create_pixel_polynomial_report(path_70, path_110)