import os
import re
import glob
import unicodedata
import numpy as np
import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
from scipy.stats import norm
from openpyxl.styles import PatternFill, Font

# ================= Settings =================
folder = r"C:\Users\dvirs\Desktop\PCB\LT\PCBA tests"  # <-- Set your folder path
EXPECTED_FRAMES = 5

# RowMap: 60 rows mapped to caps
cap1, cap2, cap3, cap4, cap5, cap6 = 1, 5, 10, 15, 20, 30
RowMap = np.zeros(60, dtype=int)
for i in range(20):
    RowMap[i] = cap1 if i % 2 == 0 else cap6
for i in range(20, 40):
    RowMap[i] = cap2 if i % 2 == 0 else cap5
for i in range(40, 60):
    RowMap[i] = cap3 if i % 2 == 0 else cap4
CAP_LIST = [cap1, cap2, cap3, cap4, cap5, cap6]
EDGE_CAP = 1
MAIN_CAPS = [c for c in CAP_LIST if c != EDGE_CAP]

# Colors
COLOR_PASS = "D1F2D1"
COLOR_WARN = "FFF4CC"
COLOR_FAIL = "F8D7DA"

# ================= Helpers =================
def pcba_sort_key(path: str):
    """
    מחזיר את הסיריאל למיון לפי סדר ייצור
    """
    name = strip_combining(os.path.basename(path))
    m = re.match(r'^(LT[_ ]?PCBA_\d{9})_(\d{8})_(\d{6})\.xlsx$', name, flags=re.IGNORECASE)
    if not m:
        return 0  # במקרה של קובץ לא תקין

    serial = int(m.group(1).split('_')[-1])
    return serial


# לדוגמה, אם רוצים את ה-LT_PCBA_XXXXXXXXX
def get_pcba_base_name(path: str) -> str:
    name = strip_combining(os.path.basename(path))
    m = re.match(r'^(LT[_ ]?PCBA_\d{9})_(\d{8})_(\d{6})\.xlsx$', name, flags=re.IGNORECASE)
    if not m:
        return ""
    return m.group(1)


def strip_combining(s: str) -> str:
    nfkd = unicodedata.normalize('NFKD', s)
    return ''.join(ch for ch in nfkd if not unicodedata.combining(ch))


def is_valid_pcba_file(path: str) -> bool:
    name = os.path.basename(path)
    name = strip_combining(name)
    pattern = r'^LT[_ ]?PCBA_\d{9}_\d{8}_\d{6}\.xlsx$'
    return re.match(pattern, name, flags=re.IGNORECASE) is not None


# ================= Load Excel Frames =================
def load_frame_cube(file_path: str, expected_frames: int = EXPECTED_FRAMES) -> np.ndarray:
    # print(f"Opening file: {file_path}")
    xls = pd.ExcelFile(file_path, engine="openpyxl")
    frame_sheets = [s for s in xls.sheet_names if s.lower().startswith("frame_")]
    if not frame_sheets:
        raise ValueError(f"No Frame_* sheets found in {os.path.basename(file_path)}.")
    try:
        frame_sheets = sorted(frame_sheets, key=lambda s: int(s.split("_")[1]))
    except Exception:
        frame_sheets = sorted(frame_sheets)
    frames = []
    for s in frame_sheets[:expected_frames]:
        df = pd.read_excel(file_path, sheet_name=s, header=None, engine="openpyxl")
        df_num = df.apply(pd.to_numeric, errors="coerce").dropna(axis=0, how="all").dropna(axis=1, how="all")
        if df_num.shape[0] < 60 or df_num.shape[1] < 30:
            raise ValueError(f"{os.path.basename(file_path)} sheet {s}: shape {df_num.shape}, expected at least (60,30).")
        frame = df_num.iloc[:60, :30].to_numpy()
        frames.append(frame)
    data = np.stack(frames, axis=2)
    if data.shape[2] != expected_frames:
        print(f"Warning: {os.path.basename(file_path)} has {data.shape[2]} frames, expected {expected_frames}.")
    return data

# ================= Analyze File =================
def analyze_file(file_path: str) -> pd.DataFrame:
    data = load_frame_cube(file_path, expected_frames=EXPECTED_FRAMES)
    rows = []

    for cap in CAP_LIST:
        mask_rows = (RowMap == cap)
        frames_masked = data[mask_rows, :, :]  # shape: (rows, cols, frames)

        # ----------------------------------------
        # Frame-to-Frame Noise – על כל הפריימים
        # ----------------------------------------
        std_per_pixel = np.std(frames_masked, axis=2)
        FrameNoise = float(np.mean(std_per_pixel))

        # ----------------------------------------
        # ממוצע על הפריימים לכל פיקסל
        # ----------------------------------------
        mean_pixels = np.mean(frames_masked, axis=2)  # shape: (rows, cols)

        # Error % – ממוצע הפיקסלים
        Error = float(np.mean(mean_pixels - cap) / cap * 100)

        # Outliers % – ממוצע הפיקסלים
        lower, upper = cap * 0.9, cap * 1.1
        total_pixels = mean_pixels.size
        outlier_count = int(np.sum((mean_pixels < lower) | (mean_pixels > upper)))
        Outliers = outlier_count / total_pixels * 100

        # Raw Outliers – כל הפריימים (1500 פיקסלים)
        all_pixels = frames_masked.flatten()
        raw_outlier_count = int(np.sum((all_pixels < lower) | (all_pixels > upper)))
        Raw_Outliers = raw_outlier_count / all_pixels.size * 100

        # Thresholds
        ERROR_FAIL = 15  # % difference
        NOISE_FAIL = 1.5  # Std deviation
        OUTLIERS_FAIL = 10  # % of pixels outside range

        # Initial status
        Status = "PASS"

        # בדיקת מדדים
        if abs(Error) > ERROR_FAIL or FrameNoise > NOISE_FAIL or Outliers > OUTLIERS_FAIL:
            # חריגות בקיבולים קיצוניים: 1pF או 30pF -> WARN במקום FAIL
            if cap in [1, 30]:
                Status = "WARN"
            else:
                Status = "FAIL"

        rows.append({
            "Cap": cap,
            "Error %": round(Error, 2),
            "Frame-to-Frame Noise (Avg)": round(FrameNoise, 2),
            "Outliers %": round(Outliers, 2),
            "Outlier Pixels": outlier_count,
            "Total Pixels": total_pixels,
            "Min Value": round(float(np.min(mean_pixels)), 2),
            "Max Value": round(float(np.max(mean_pixels)), 2),
            "Status": Status,
            "Raw Outliers %": round(Raw_Outliers, 2)
        })

    df_result = pd.DataFrame(rows)
    df_result["File"] = os.path.basename(file_path)
    return df_result, data


# ================= Collect Files =================
patterns = ["*.xlsx"]
all_files = []
for pat in patterns:
    all_files.extend(glob.glob(os.path.join(folder, pat)))
all_files = [f for f in all_files if is_valid_pcba_file(f)]
all_files = sorted(all_files, key=pcba_sort_key)

if not all_files:
    print("No Excel files found.")
    raise SystemExit(1)

# ================= Run Analysis =================
final_report = []
errors = []
all_data = {}  # store all loaded data per file
total = len(all_files)
print(f"Found {total} files.\n")

for i, f in enumerate(all_files, start=1):
    fname = os.path.basename(f)
    print(f"[{i}/{total}] Opening: {fname} ...")
    try:
        df, data = analyze_file(f)
        final_report.append(df)
        all_data[fname] = data
        print(f"[{i}/{total}] ✔ OK: {fname}, added {len(df)} rows.\n")
    except Exception as e:
        errors.append((fname, str(e)))
        print(f"[{i}/{total}] ✖ ERROR in {fname}: {e}\n")

report_df = pd.concat(final_report, ignore_index=True) if final_report else pd.DataFrame()

# ================= Boards Summary =================
summary_rows = []

for board, grp in report_df.groupby("File"):

    cap_status_map = dict(zip(grp["Cap"], grp["Status"]))

    # הכללות
    fail_caps = [cap for cap, st in cap_status_map.items() if st == "FAIL" and cap not in [1, 30]]
    warn_caps = [cap for cap, st in cap_status_map.items() if st == "WARN" or (st == "FAIL" and cap in [1, 30])]

    # מצב לוח
    overall_status = "FAIL" if fail_caps else "PASS"

    # overall_score = float(grp["Score"].mean())

    summary_rows.append({
        "File": board,
        # "Overall Score": round(overall_score, 2),
        "Overall Status": overall_status,
        "Fail Caps": ", ".join(str(c) for c in fail_caps) if fail_caps else "",
        "Warn Caps": ", ".join(str(c) for c in warn_caps) if warn_caps else "",
        "PASS Count": int((grp["Status"] == "PASS").sum())
    })

boards_summary = pd.DataFrame(summary_rows)


excel_path = os.path.join(folder, "Final_Report.xlsx")
with pd.ExcelWriter(excel_path, engine="openpyxl") as writer:
    report_df.to_excel(writer, index=False, sheet_name="Per_Cap")
    boards_summary.to_excel(writer, index=False, sheet_name="Boards_Summary")
print(f"✅ Excel saved: {excel_path}")

# Function to update the title size on each slide
def set_title_size(slide, size_pt=24):
    if slide.shapes.title:
        title_tf = slide.shapes.title.text_frame
        for paragraph in title_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size_pt)
                run.font.bold = True  # אפשר להשאיר כ-bold
                run.font.color.rgb = RGBColor(0,0,0)

# ================= PowerPoint =================
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "PCBA Measurement Quality Report"
slide.placeholders[1].text = f"FW / HW Version: TBD\nDate: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

# ================= Parameters Slide =================
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "PCBA Measurement Metrics – Logic & Interpretation"

text = (
    "🔹 Frame-to-Frame Noise (Avg):\n"
    "    - Measures the variation of each pixel across 5 frames.\n"
    "    - Reflects momentary noise in the measurement / ADC.\n"
    "    - Calculated using all pixels across all frames.\n\n"

    "🔹 Error %:\n"
    "    - Accuracy: percentage difference between the mean pixel value and expected capacitance.\n"
    "    - Calculated after averaging each pixel across 5 frames.\n"
    "    - Low value = accurate measurement.\n\n"

    "🔹 Outliers %:\n"
    "    - Percentage of pixels outside ±10% of expected capacitance.\n"
    "    - Calculated on the per-pixel averages.\n"
    "    - High value = unstable or inconsistent measurement.\n\n"

    "🔹 Outlier Pixels:\n"
    "    - Absolute number of outlier pixels.\n\n"

    "🔹 Total Pixels:\n"
    "    - Total number of pixels measured per capacitance.\n\n"

    "🔹 Min / Max Value:\n"
    "    - Minimum and maximum of the per-pixel averages.\n"
    "    - Gives insight into the measurement range.\n\n"

    "🔹 Status (PASS / WARN / FAIL):\n"
    "    - Determines the overall quality per capacitance.\n"
    "    - FAIL → one or more non-edge capacitors exceed thresholds (Error %, Frame-to-Frame Noise, or Outliers %).\n"
    "    - WARN → edge capacitors (1pF or 30pF) exceed thresholds; considered a caution, not a failure.\n"
    "    - PASS → all metrics within thresholds, or only edge capacitors triggered WARN.\n"
    "    - In summary charts (pie charts), WARN is grouped with PASS for simplicity."

)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(5))
tf = txBox.text_frame
tf.clear()
p = tf.add_paragraph()
p.text = text
p.font.size = Pt(12)
set_title_size(slide, 24)

# ================= Per-Board Slides =================
boards = boards_summary["File"].tolist()
offset_down = Inches(0.6)
for board in boards:
    df_board = report_df[report_df["File"]==board].copy()
    df_cap = df_board.groupby('Cap', as_index=False).agg({
        'Frame-to-Frame Noise (Avg)': 'mean',
        'Outliers %': 'mean',
        'Outlier Pixels': 'sum',
        'Total Pixels': 'first',
        'Min Value': 'min',
        'Max Value': 'max',
        # 'Score': 'mean',
        'Error %': 'mean',
        'Status': lambda s: (
            'FAIL' if (s == 'FAIL').any()
            else 'WARN' if (s == 'WARN').any()
            else 'PASS'
        )
    })

    df_main = df_cap[df_cap["Cap"].isin(MAIN_CAPS)]
    df_edge = df_cap[df_cap["Cap"] == EDGE_CAP]

    fail_caps = df_cap[df_cap["Status"] == "FAIL"]["Cap"].tolist()
    warn_caps = df_cap[df_cap["Status"] == "WARN"]["Cap"].tolist()

    fail_count = len(fail_caps)
    warn_count = len(warn_caps)

    if fail_count >= 2:
        overall_status = "FAIL"
    elif fail_count == 1:
        overall_status = "WARN"
    elif warn_count >= 2:
        overall_status = "WARN"
    else:
        overall_status = "PASS"

    edge_status = (
        df_edge["Status"].iloc[0]
        if not df_edge.empty
        else "N/A"
    )

    board_name = get_pcba_base_name(board)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Board: {board_name}"
    title_tf = slide.shapes.title.text_frame
    title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    set_title_size(slide, 24)

    # Table
    rows_table = len(CAP_LIST)+1
    cols_table = 10
    headers = [
        'Cap',
        # 'Score',
        'Error %',
        'Frame-to-Frame Noise (Avg)',
        'Outliers %',
        'Outlier Pixels',
        'Total Pixels',
        'Min',
        'Max',
        'Status'
    ]

    table = slide.shapes.add_table(
        rows_table,
        cols_table,
        Inches(0.3),
        Inches(1) + offset_down,
        Inches(9.5),
        Inches(2.2)
    ).table
    for c, h in enumerate(headers):
        table.cell(0,c).text = h

    std_vals, out_vals = [], []
    for r, cap in enumerate(CAP_LIST):
        row = df_cap[df_cap['Cap']==cap]
        table.cell(r+1,0).text = str(cap)
        if row.empty:
            table.cell(r+1,1).text = "-"
            table.cell(r+1,2).text = "-"
            table.cell(r+1,3).text = "-"
            table.cell(r+1,4).text = "-"
            table.cell(r+1,5).text = "-"
            std_vals.append(0.0)
            out_vals.append(0.0)
        else:
            row_data = row.iloc[0]
            # table.cell(r + 1, 1).text = f"{row_data['Score']:.2f}"
            table.cell(r + 1, 2).text = f"{row_data['Error %']:.2f}"
            table.cell(r + 1, 3).text = f"{row_data['Frame-to-Frame Noise (Avg)']:.2f}"
            table.cell(r + 1, 4).text = f"{row_data['Outliers %']:.2f}"

            table.cell(r + 1, 5).text = str(int(row_data['Outlier Pixels']))
            table.cell(r + 1, 6).text = str(int(row_data['Total Pixels']))
            table.cell(r + 1, 7).text = f"{row_data['Min Value']:.2f}"
            table.cell(r + 1, 8).text = f"{row_data['Max Value']:.2f}"

            table.cell(r + 1, 9).text = row_data['Status']
            # Color Status cell
            cell = table.cell(r+1,9)
            cell.fill.solid()
            if row_data['Status']=="PASS":
                cell.fill.fore_color.rgb = RGBColor(0xD1,0xF2,0xD1)
            elif row_data['Status']=="WARN":
                cell.fill.fore_color.rgb = RGBColor(0xFF,0xF4,0xCC)
            else:
                cell.fill.fore_color.rgb = RGBColor(0xF8,0xD7,0xDA)
            std_vals.append(float(row_data['Frame-to-Frame Noise (Avg)']))
            out_vals.append(float(row_data['Outliers %']))

    # Overall status box
    left, top, width, height = Inches(5.7), Inches(0.1) + offset_down, Inches(4), Inches(0.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        f"Overall Status: {overall_status}\n"
        f"Fail Caps: {', '.join(map(str, fail_caps)) if fail_caps else 'None'}\n"
        f"Warn Caps: {', '.join(map(str, warn_caps)) if warn_caps else 'None'}"
    )

    run.font.bold = True
    run.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    fill = box.fill
    fill.solid()
    if overall_status=="PASS":
        fill.fore_color.rgb = RGBColor(0xD1,0xF2,0xD1)
    elif overall_status=="WARN":
        fill.fore_color.rgb = RGBColor(0xFF,0xF4,0xCC)
    else:
        fill.fore_color.rgb = RGBColor(0xF8,0xD7,0xDA)
    for run in p.runs:
        run.font.color.rgb = RGBColor(0,0,0)

    # Charts
    std_path = os.path.join(folder, f"{board}_Noise.png")
    out_path = os.path.join(folder, f"{board}_Outliers.png")
    fig, ax = plt.subplots(figsize=(6,3))
    ax.bar(CAP_LIST, std_vals, color='skyblue')
    ax.set_title("Noise vs Capacitance")
    ax.set_xlabel("Capacitance (pF)")
    ax.set_ylabel("Noise (pF)")
    plt.tight_layout()
    plt.savefig(std_path)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(6,3))
    ax.bar(CAP_LIST, out_vals, color='salmon')
    ax.set_title("Outliers % vs Capacitance")
    ax.set_xlabel("Capacitance (pF)")
    ax.set_ylabel("Outliers %")
    plt.tight_layout()
    plt.savefig(out_path)
    plt.close(fig)

    slide.shapes.add_picture(std_path, Inches(0.3), Inches(5) + offset_down, width=Inches(3.6), height=Inches(1.8))
    slide.shapes.add_picture(out_path, Inches(4), Inches(5) + offset_down, width=Inches(3.6), height=Inches(1.8))
    for tmp in (std_path, out_path):
        try: os.remove(tmp)
        except: pass

# ================= Per-Cap Slides =================
for cap in CAP_LIST:
    all_pixels = []
    avg_per_board = []

    for board, data in all_data.items():
        mask_rows = (RowMap == cap)
        frames_masked = data[mask_rows, :, :]

        all_pixels.extend(frames_masked.flatten())

        mean_board = float(np.mean(frames_masked))  # ממוצע כולל לכל הלוח
        avg_per_board.append(mean_board)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Cap {cap}pF: All Pixels + Board Means"
    set_title_size(slide, 24)

    # =========================
    # Histogram – all pixels
    # =========================
    fig, ax = plt.subplots(figsize=(6, 3))
    ax.hist(all_pixels, bins=50, edgecolor='black')
    ax.set_title(f"All measurements for {cap}pF")
    ax.set_xlabel("Pixel value")
    ax.set_ylabel("Frequency")
    plt.tight_layout()

    path_all = os.path.join(folder, f"{cap}_all_pixels.png")
    plt.savefig(path_all)
    plt.close(fig)

    slide.shapes.add_picture(
        path_all,
        Inches(0.5),
        Inches(1.2),
        width=Inches(5),
        height=Inches(1.8)
    )

    try:
        os.remove(path_all)
    except:
        pass

    # =========================
    # Histogram – average per board
    # =========================
    fig, ax = plt.subplots(figsize=(6, 3))
    ax.hist(avg_per_board, bins=50, edgecolor='black')
    ax.set_title(f"Average per board for {cap}pF")
    ax.set_xlabel("Mean Pixel Value per Board")
    ax.set_ylabel("Frequency")
    plt.tight_layout()

    path_avg = os.path.join(folder, f"{cap}_avg_per_board.png")
    plt.savefig(path_avg)
    plt.close(fig)

    slide.shapes.add_picture(
        path_avg,
        Inches(0.5),
        Inches(3.2),
        width=Inches(5),
        height=Inches(1.8)
    )

    try:
        os.remove(path_avg)
    except:
        pass

    # =========================
    # Gaussian distribution + σ lines
    # =========================
    mu = np.mean(all_pixels)
    sigma = np.std(all_pixels)

    fig, ax = plt.subplots(figsize=(6, 3))

    # Normalized histogram
    ax.hist(
        all_pixels,
        bins=50,
        density=True,
        alpha=0.6,
        edgecolor='black'
    )

    # Gaussian curve
    x = np.linspace(min(all_pixels), max(all_pixels), 500)
    pdf = norm.pdf(x, mu, sigma)
    ax.plot(x, pdf, linewidth=2)

    # Std deviation markers
    for k in [1, 2, 3]:
        ax.axvline(mu + k * sigma, linestyle='--', linewidth=1)
        ax.axvline(mu - k * sigma, linestyle='--', linewidth=1)

    ax.set_title(
        f"Gaussian Distribution for {cap}pF\n"
        f"μ={mu:.2f}, σ={sigma:.2f}"
    )
    ax.set_xlabel("Pixel value")
    ax.set_ylabel("Probability Density")

    plt.tight_layout()

    path_gauss = os.path.join(folder, f"{cap}_gaussian.png")
    plt.savefig(path_gauss)
    plt.close(fig)

    slide.shapes.add_picture(
        path_gauss,
        Inches(0.5),
        Inches(5.4),
        width=Inches(5),
        height=Inches(1.8)
    )

    try:
        os.remove(path_gauss)
    except:
        pass

# ================= Per-Cap Status Pie Charts + Summary Table =================
for cap in CAP_LIST:

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Cap {cap}pF – Status Distribution + Summary"
    set_title_size(slide, 24)

    # =========================
    # Aggregate all pixels across boards
    # =========================
    all_pixels = []
    for board, data in all_data.items():
        mask_rows = (RowMap == cap)
        frames_masked = data[mask_rows, :, :]
        all_pixels.extend(frames_masked.flatten())
    all_pixels = np.array(all_pixels)

    mean_val = np.mean(all_pixels)
    min_val = np.min(all_pixels)
    max_val = np.max(all_pixels)
    std_val = np.std(all_pixels)
    lower, upper = cap*0.9, cap*1.1
    outlier_pixels = int(np.sum((all_pixels < lower) | (all_pixels > upper)))
    total_pixels = all_pixels.size
    outlier_percent = outlier_pixels / total_pixels * 100

    # =========================
    # Pie chart – PASS/WARN/FAIL
    # =========================
    df_cap = report_df[report_df["Cap"] == cap]
    status_counts = df_cap["Status"].value_counts()
    labels = status_counts.index.tolist()
    sizes = status_counts.values.tolist()
    colors_hex = {"PASS": "#D1F2D1", "WARN": "#FFF4CC", "FAIL": "#F8D7DA"}
    colors = [colors_hex.get(l, "#CCCCCC") for l in labels]

    fig, ax = plt.subplots(figsize=(4,4))
    ax.pie(sizes, labels=labels, autopct="%1.1f%%", startangle=90, colors=colors)
    ax.set_title(f"{cap}pF – PASS / WARN / FAIL")
    plt.tight_layout()
    pie_path = os.path.join(folder, f"Pie_{cap}pF.png")
    plt.savefig(pie_path)
    plt.close(fig)

    # הצגת הפאי בסלייד
    slide.shapes.add_picture(
        pie_path,
        Inches(2),
        Inches(1.1),
        width=Inches(5),
        height=Inches(5)
    )
    try:
        os.remove(pie_path)
    except:
        pass

    # =========================
    # Summary table of all pixels
    # =========================
    rows_table = 2
    cols_table = 8
    headers = ["Cap", "Mean", "Min", "Max", "Std Dev", "Outliers %", "Outlier Pixels", "Total Pixels"]

    table = slide.shapes.add_table(
        rows_table, cols_table,
        Inches(0.5), Inches(6),  # מיקום הטבלה
        Inches(8), Inches(1)        # גודל הטבלה
    ).table

    # כותרות
    for c, h in enumerate(headers):
        table.cell(0, c).text = h

    # שורה עם הערכים
    table.cell(1, 0).text = str(cap)
    table.cell(1, 1).text = f"{mean_val:.2f}"
    table.cell(1, 2).text = f"{min_val:.2f}"
    table.cell(1, 3).text = f"{max_val:.2f}"
    table.cell(1, 4).text = f"{std_val:.2f}"
    table.cell(1, 5).text = f"{outlier_percent:.2f}"
    table.cell(1, 6).text = str(outlier_pixels)
    table.cell(1, 7).text = str(total_pixels)



# ================= Save PPT =================
ppt_path = os.path.join(folder,"PCBA_Report_Full.pptx")
prs.save(ppt_path)
print(f"🖼️ PowerPoint saved: {ppt_path}")
print("Done.")
